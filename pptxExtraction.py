from pptx import Presentation
from pathlib import Path

ppt_path = Path("/Users/leichunyin/Documents/In Birmingham/During Study/Modules/2Genomics&NGS/Week1/Hidden Markov Model/HMM_lecture.pptx")
output_path = "ppt_text.txt"
prs = Presentation(ppt_path)

with open(output_path, "w", encoding="utf-8") as f:
    for i, slide in enumerate(prs.slides, start=1):
        f.write(f"\n--- Slide {i} ---\n")

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    f.write(text + "\n")
